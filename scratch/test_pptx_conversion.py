"""Quick smoke test for PPTX -> PDF conversion."""
import asyncio, io, sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import fitz
from pptx import Presentation
from app.utils.document_converter import convert_to_pdf

async def main():
    # Create a test PPTX in memory
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Slide Title"
    slide.placeholders[1].text = "Body content for verification"

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    pptx_bytes = buf.read()
    print(f"PPTX size: {len(pptx_bytes)} bytes")

    # Convert
    pdf_bytes = await convert_to_pdf(pptx_bytes, "test.pptx")
    if pdf_bytes is None:
        print("CONVERSION FAILED")
        return

    print(f"PDF size: {len(pdf_bytes)} bytes")

    # Verify PDF is readable
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    print(f"PDF pages: {doc.page_count}")
    text = doc[0].get_text()
    print(f"Page 1 text: {repr(text[:200])}")
    doc.close()
    print("SUCCESS")

asyncio.run(main())
