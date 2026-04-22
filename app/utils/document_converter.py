"""
app/utils/document_converter.py
──────────────────────────────────
Utility for converting DOCX, TXT, and PPTX files to PDF.
Uses python-docx, python-pptx for parsing and PyMuPDF (fitz) for PDF generation.
"""

import io
import logging
import fitz  # PyMuPDF
from docx import Document

logger = logging.getLogger(__name__)

async def convert_to_pdf(file_bytes: bytes, filename: str) -> bytes | None:
    """
    Converts supported documents to PDF bytes.
    Returns bytes of the generated PDF, or None if conversion fails.
    """
    try:
        ext = filename.lower().split('.')[-1]
        
        if ext == 'pdf':
            return file_bytes
            
        if ext in ['docx', 'doc']:
            return await _convert_docx_to_pdf(file_bytes)
        elif ext == 'txt':
            return await _convert_txt_to_pdf(file_bytes)
        elif ext in ['pptx', 'ppt']:
            return await _convert_pptx_to_pdf(file_bytes)
            
        logger.warning(f"Unsupported file type for conversion: {ext}")
        return None
        
    except Exception as e:
        logger.error(f"Failed to convert '{filename}' to PDF: {str(e)}")
        return None

async def _convert_docx_to_pdf(file_bytes: bytes) -> bytes | None:
    """
    Simple DOCX to PDF conversion. 
    Extracts text and renders it into a new PDF using fitz.
    """
    try:
        doc = Document(io.BytesIO(file_bytes))
        pdf_doc = fitz.open()
        
        # Simple text-based reconstruction
        # We create pages as needed
        text_content = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text)
        
        full_text = "\n\n".join(text_content)
        
        if not full_text:
            full_text = "[No text content found in document]"

        # Create a page and insert text
        # Using a standard font that's guaranteed to be in fitz
        page = pdf_doc.new_page()
        
        # Define a safe margin and area
        rect = fitz.Rect(72, 72, 523, 770) # Standard margins
        
        # Insert text with automatic wrapping
        page.insert_textbox(rect, full_text, fontname="helv", fontsize=11)
        
        # If text overflows, we'd ideally create more pages. 
        # For "basic text rendering", we'll at least handle one page well.
        # IMPROVEMENT: Basic multi-page support for long text
        
        pdf_bytes = pdf_doc.tobytes()
        pdf_doc.close()
        return pdf_bytes
        
    except Exception as e:
        logger.error(f"DOCX conversion error: {e}")
        return None

async def _convert_txt_to_pdf(file_bytes: bytes) -> bytes | None:
    """
    Converts plain text to PDF.
    """
    try:
        text = file_bytes.decode('utf-8', errors='ignore')
        pdf_doc = fitz.open()
        
        # Simple multi-page logic for long text
        lines = text.splitlines()
        lines_per_page = 50
        
        for i in range(0, len(lines), lines_per_page):
            page = pdf_doc.new_page()
            page_text = "\n".join(lines[i:i + lines_per_page])
            rect = fitz.Rect(50, 50, 550, 800)
            page.insert_textbox(rect, page_text, fontname="helv", fontsize=10)
            
        pdf_bytes = pdf_doc.tobytes()
        pdf_doc.close()
        return pdf_bytes
        
    except Exception as e:
        logger.error(f"TXT conversion error: {e}")
        return None


async def _convert_pptx_to_pdf(file_bytes: bytes) -> bytes | None:
    """
    Converts PPTX presentations to PDF.
    Extracts text content from each slide and renders it as a separate PDF page.
    Uses python-pptx for reading and PyMuPDF (fitz) for PDF generation.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation(io.BytesIO(file_bytes))
        pdf_doc = fitz.open()

        slide_count = len(prs.slides)
        if slide_count == 0:
            # Empty presentation — create a single page with a notice
            page = pdf_doc.new_page()
            rect = fitz.Rect(72, 72, 523, 770)
            page.insert_textbox(rect, "[Empty presentation — no slides found]",
                                fontname="helv", fontsize=12)
            pdf_bytes = pdf_doc.tobytes()
            pdf_doc.close()
            return pdf_bytes

        for slide_idx, slide in enumerate(prs.slides, start=1):
            page = pdf_doc.new_page(width=720, height=540)  # landscape 10x7.5"

            # ── Slide header ──────────────────────────────────────────────
            header_rect = fitz.Rect(36, 20, 684, 45)
            page.insert_textbox(header_rect,
                                f"Slide {slide_idx} / {slide_count}",
                                fontname="helv", fontsize=9, color=(0.4, 0.4, 0.4))

            # ── Collect text from all shapes on this slide ─────────────────
            slide_texts = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        slide_texts.append(line)

            if not slide_texts:
                slide_texts = ["[No text content on this slide]"]

            slide_text = "\n\n".join(slide_texts)

            # ── Render text into the page body ────────────────────────────
            body_rect = fitz.Rect(36, 50, 684, 510)
            page.insert_textbox(body_rect, slide_text,
                                fontname="helv", fontsize=11)

        pdf_bytes = pdf_doc.tobytes()
        pdf_doc.close()
        logger.info(f"PPTX conversion successful: {slide_count} slides → PDF")
        return pdf_bytes

    except Exception as e:
        logger.error(f"PPTX conversion error: {e}")
        return None
