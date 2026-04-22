"""
app/utils/text_extractor.py
────────────────────────────
Shared utility for extracting plain text from PDF and DOCX files.
Used by both ai_service and chatbot_service.
"""

import logging

import logging

logger = logging.getLogger(__name__)


async def extract_text(file_path: str) -> list[dict]:
    import os
    import re
    import io

    # Detect if it's an S3 public URL
    is_url = file_path.startswith("http://") or file_path.startswith("https://")
    
    file_bytes = None

    if is_url:
        import httpx
        try:
            async with httpx.AsyncClient(timeout=10.0) as client:
                response = await client.get(file_path)
                response.raise_for_status()
                file_bytes = response.content
        except Exception as e:
            logger.warning(f"Failed to fetch document from URL '{file_path}': {e}")
            return []
    else:
        file_path = os.path.normpath(file_path)
        if not os.path.isabs(file_path):
            file_path = os.path.abspath(file_path)
        
        if not os.path.exists(file_path) and "uploads" in file_path:
            # Fallback relative logic
            rel_path = file_path.split("uploads")[-1].lstrip("\\/")
            file_path = os.path.abspath(os.path.join("uploads", rel_path))
            
        if not os.path.exists(file_path):
            logger.warning(f"File not found for extraction: {file_path}")
            return []
            
        with open(file_path, "rb") as f:
            file_bytes = f.read()

    def clean_text(text: str) -> str:
        text = re.sub(r'\s+', ' ', text)
        return text.strip()

    def _sync_extract():
        try:
            if file_path.lower().split("?")[0].endswith(".pdf"):
                import fitz  # PyMuPDF
                pages: list[dict] = []
                
                with fitz.open(stream=file_bytes, filetype="pdf") as doc:
                    for i, page in enumerate(doc):
                        text = clean_text(page.get_text())
                        if text:
                            pages.append({"text": text, "page": i + 1})
                return pages

            elif file_path.lower().split("?")[0].endswith((".docx", ".doc")):
                from docx import Document
                doc = Document(io.BytesIO(file_bytes))
                full_text = clean_text("\n".join(p.text for p in doc.paragraphs if p.text.strip()))
                if full_text:
                    return [{"text": full_text, "page": 1}]
                return []

            elif file_path.lower().split("?")[0].endswith(".txt"):
                text = clean_text(file_bytes.decode('utf-8', errors="ignore"))
                if text:
                    return [{"text": text, "page": 1}]
                return []

            elif file_path.lower().split("?")[0].endswith((".pptx", ".ppt")):
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_bytes))
                pages = []
                for i, slide in enumerate(prs.slides):
                    slide_texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                line = para.text.strip()
                                if line:
                                    slide_texts.append(line)
                    if slide_texts:
                        pages.append({"text": clean_text(" ".join(slide_texts)), "page": i + 1})
                return pages

        except Exception as e:
            logger.warning(f"Internal extraction logic failed for '{file_path}': {e}")
            return []
        return []

    import asyncio
    return await asyncio.to_thread(_sync_extract)
